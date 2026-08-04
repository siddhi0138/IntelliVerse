"""Document intelligence: the actual reason Qdrant earns a place in this
stack. Every other "knowledge base" in this app is already structured
JSON (schema, stats, findings) — a vector database over that would be
adopting a technology nothing here needs. Unstructured documents (PDFs,
reports, SOPs, meeting notes) are different: there IS something to
retrieve, so retrieval-then-narrate is a genuine, not performative, use
of the pattern every other LLM call in this backend already follows.

Embeddings run locally via fastembed — no external API key, consistent
with the rest of this project's "self-hosted over SaaS" bias
(Prometheus/Grafana over Sentry/Langfuse). fastembed (not
sentence-transformers) specifically: it runs the same MiniLM-family model
through ONNX Runtime instead of PyTorch, which avoids the ~500MB+ of
resident memory PyTorch alone costs — on a memory-constrained deployment
(e.g. Render's free tier, ~512MB) that difference was the gap between this
feature working and it OOM-killing the entire backend process the first
time anyone uploaded a document, not just this feature returning an error.
Qdrant runs in local on-disk mode (no separate server process) for the
same reason SQLite backs the dataset catalog: this is a local-first tool,
and a full Qdrant server is unneeded complexity until there's a reason to
scale retrieval across processes.

LlamaIndex wasn't adopted for the chunk/embed/retrieve pipeline itself:
its real value is multi-source query orchestration, which a single
document collection with one retrieval step doesn't need — direct
qdrant-client + sentence-transformers is simpler to read and debug for
what's actually being done here.
"""

from __future__ import annotations

import io
import os
import re
import uuid
from pathlib import Path
from typing import TYPE_CHECKING

import docx
import pypdf
from pptx import Presentation
from qdrant_client import QdrantClient
from qdrant_client.models import (
    Distance,
    FieldCondition,
    Filter,
    MatchValue,
    PayloadSchemaType,
    PointStruct,
    VectorParams,
)

if TYPE_CHECKING:
    # fastembed still isn't free (ONNX Runtime + the model weights), so it
    # stays deferred to _get_model() below — a process that never touches
    # document Q&A never pays for it.
    from fastembed import TextEmbedding

_QDRANT_PATH = Path(__file__).parent / "data" / "qdrant"
_QDRANT_URL = os.environ.get("QDRANT_URL")
_QDRANT_API_KEY = os.environ.get("QDRANT_API_KEY")
_COLLECTION = "documents"
_EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
_EMBEDDING_DIM = 384

_client: QdrantClient | None = None
_model: TextEmbedding | None = None


class UnsupportedDocumentError(Exception):
    pass


def _get_client() -> QdrantClient:
    global _client
    if _client is None:
        # QDRANT_URL set (e.g. a Qdrant Cloud free-tier cluster) means a
        # real hosted instance — required in production, since local
        # on-disk mode lives on the backend's own ephemeral disk and gets
        # wiped on every Render redeploy. Local dev without QDRANT_URL set
        # keeps using on-disk mode, same as always.
        if _QDRANT_URL:
            _client = QdrantClient(url=_QDRANT_URL, api_key=_QDRANT_API_KEY)
        else:
            _QDRANT_PATH.mkdir(parents=True, exist_ok=True)
            _client = QdrantClient(path=str(_QDRANT_PATH))
        if not _client.collection_exists(_COLLECTION):
            _client.create_collection(
                _COLLECTION,
                vectors_config=VectorParams(size=_EMBEDDING_DIM, distance=Distance.COSINE),
            )
        # Local on-disk mode filters on any payload field with a full scan,
        # no index required — Qdrant Cloud's server enforces an explicit
        # index for any field used in a query/delete filter and 400s
        # otherwise ("Index required but not found"). Both fields below are
        # filtered on in search()/delete_document(). Idempotent (safe to
        # call even if the index already exists), and cheap since this only
        # runs once per process (guarded by `_client is None` above).
        _client.create_payload_index(_COLLECTION, "username", field_schema=PayloadSchemaType.KEYWORD)
        _client.create_payload_index(_COLLECTION, "doc_id", field_schema=PayloadSchemaType.KEYWORD)
    return _client


def _get_model() -> "TextEmbedding":
    global _model
    if _model is None:
        from fastembed import TextEmbedding

        _model = TextEmbedding(model_name=_EMBEDDING_MODEL_NAME)
    return _model


def extract_text(filename: str, content: bytes) -> str:
    lowered = filename.lower()
    if lowered.endswith(".pdf"):
        reader = pypdf.PdfReader(io.BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    if lowered.endswith(".docx"):
        document = docx.Document(io.BytesIO(content))
        return "\n".join(p.text for p in document.paragraphs)
    if lowered.endswith(".pptx"):
        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
        return "\n".join(texts)
    if lowered.endswith(".txt"):
        return content.decode("utf-8", errors="ignore")
    raise UnsupportedDocumentError(f"Unsupported document type: {filename}. Use PDF, DOCX, PPTX, or TXT.")


def chunk_text(text: str, chunk_size: int = 800, overlap: int = 150) -> list[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        if end >= len(text):
            break
        start = end - overlap
    return chunks


def ingest_document(username: str, doc_id: str, filename: str, content: bytes) -> int:
    """Extracts, chunks, embeds, and stores a document. Returns the chunk count."""
    text = extract_text(filename, content)
    chunks = chunk_text(text)
    if not chunks:
        return 0

    model = _get_model()
    embeddings = list(model.embed(chunks))
    client = _get_client()
    points = [
        PointStruct(
            id=str(uuid.uuid4()),
            vector=embedding.tolist(),
            payload={
                "username": username,
                "doc_id": doc_id,
                "filename": filename,
                "chunk_index": i,
                "text": chunk,
            },
        )
        for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
    ]
    client.upsert(_COLLECTION, points)
    return len(chunks)


def search(username: str, query: str, limit: int = 5) -> list[dict]:
    client = _get_client()
    model = _get_model()
    query_vector = next(model.embed([query])).tolist()
    results = client.query_points(
        _COLLECTION,
        query=query_vector,
        query_filter=Filter(must=[FieldCondition(key="username", match=MatchValue(value=username))]),
        limit=limit,
    ).points
    return [
        {
            "filename": p.payload["filename"],
            "chunk_index": p.payload["chunk_index"],
            "text": p.payload["text"],
            "score": round(p.score, 4),
        }
        for p in results
    ]


def delete_document(username: str, doc_id: str) -> None:
    client = _get_client()
    client.delete(
        _COLLECTION,
        points_selector=Filter(
            must=[
                FieldCondition(key="username", match=MatchValue(value=username)),
                FieldCondition(key="doc_id", match=MatchValue(value=doc_id)),
            ]
        ),
    )
