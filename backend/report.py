"""Additive: export a completed analysis as Excel, PDF, or PowerPoint.

All three builders read from the same `result` dict returned by
POST /api/analyze (plus the original DataFrame for the Excel raw-data
sheet) — no re-computation, just formatting what's already been derived.

Findings and risk alerts are shown two ways: a plain-English sentence a
non-technical reader can act on immediately, plus the underlying stat
(r value, p value, Cramer's V, confidence %, ...) alongside it — so the
numbers a technical reader wants to check are never hidden, just not
the first thing you have to parse.

The PDF and PPTX builders cover every deterministic section the app
computes (business health, data quality detail, root cause, relationships,
segmentation, schema/data dictionary) — not just findings and forecast —
so the exported report stands on its own without the app open.
"""
from __future__ import annotations

import io
import re
from datetime import datetime, timezone
from typing import Any

import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import HRFlowable, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

MAX_EXCEL_DATA_ROWS = 5000

BRAND = RGBColor(0x4F, 0x46, 0xE5)  # indigo-600, matches the app's UI
SUCCESS = RGBColor(0x05, 0x96, 0x69)  # emerald-600
CRITICAL = RGBColor(0xDC, 0x26, 0x26)  # red-600
WARNING = RGBColor(0xD9, 0x77, 0x06)  # amber-600
MUTED = RGBColor(0x64, 0x74, 0x8B)  # slate-500
LIGHT_ROW = RGBColor(0xF8, 0xFA, 0xFC)  # slate-50

BRAND_HEX = "#4f46e5"
SUCCESS_HEX = "#059669"
CRITICAL_HEX = "#dc2626"
WARNING_HEX = "#d97706"
MUTED_HEX = "#64748b"

_TRAILING_STATS = re.compile(r"\s*\([^()]*\)\s*$")

# Same thresholds as BusinessHealthPanel.tsx, so the report and the app never disagree.
_COMPONENT_LABELS = {
    "data_quality": "Data Quality",
    "growth": "Growth",
    "forecast_reliability": "Forecast Confidence",
    "safety": "Risk Safety",
}


def _health_color(score: float) -> RGBColor:
    if score >= 70:
        return SUCCESS
    if score >= 45:
        return WARNING
    return CRITICAL


def _health_color_hex(score: float) -> str:
    if score >= 70:
        return SUCCESS_HEX
    if score >= 45:
        return WARNING_HEX
    return CRITICAL_HEX


def _plain_headline(headline: str) -> str:
    """Strips a trailing "(r=0.82)" / "(Cramer's V=0.41)" / "(anova, p=0.03)"
    parenthetical so a non-technical reader gets the plain-English sentence
    first. The full headline (with the stat) is always shown alongside this,
    never dropped — just not the first thing you have to parse."""
    plain = _TRAILING_STATS.sub("", headline).strip()
    return plain or headline


def _format_alert(alert: dict[str, Any]) -> tuple[str, str]:
    """Mirrors RiskAlertsPanel.tsx's phrasing so the report and the app never disagree."""
    confidence = f" (confidence: {alert['confidence_pct']}%)" if alert.get("confidence_pct") is not None else ""
    if alert.get("kind") == "threshold_crossing":
        message = (
            f"{alert.get('metric', 'This metric')} is projected to reach a critical level (zero) within "
            f"{alert.get('periods_until_critical', '?')} period(s){confidence}."
        )
        return "Critical", message
    driver = f" Primary historical driver: {alert['primary_driver']}." if alert.get("primary_driver") else ""
    message = f"{alert.get('metric', 'This metric')} is projected to {alert.get('direction', 'change')}{confidence}.{driver}"
    return "Warning", message


def build_excel_report(result: dict[str, Any], df: pd.DataFrame) -> bytes:
    wb = Workbook()

    summary = wb.active
    summary.title = "Summary"
    summary.append(["Filename", result.get("filename")])
    summary.append(["Domain", result.get("domain")])
    summary.append(["Rows", result.get("row_count")])
    summary.append(["Columns", result.get("column_count")])
    quality = result.get("quality") or {}
    summary.append(["Data quality score", quality.get("score")])
    health = result.get("business_health") or {}
    summary.append(["Business health score", health.get("overall")])
    for row in summary.iter_rows(min_row=1, max_row=summary.max_row, min_col=1, max_col=1):
        row[0].font = Font(bold=True)
    summary.column_dimensions["A"].width = 22
    summary.column_dimensions["B"].width = 40

    schema_ws = wb.create_sheet("Schema")
    schema_ws.append(["Column", "Type", "Semantic label", "Confidence"])
    for cell in schema_ws[1]:
        cell.font = Font(bold=True)
    for col in result.get("schema", []):
        schema_ws.append([col.get("name"), col.get("type"), col.get("semantic_label"), col.get("confidence")])

    findings_ws = wb.create_sheet("Findings")
    findings_ws.append(["Rank", "Finding", "Statistical detail", "Score"])
    for cell in findings_ws[1]:
        cell.font = Font(bold=True)
    for i, f in enumerate(result.get("ranked_findings", []), start=1):
        headline = f.get("headline", "")
        findings_ws.append([i, _plain_headline(headline), headline, f.get("score")])
    findings_ws.column_dimensions["B"].width = 45
    findings_ws.column_dimensions["C"].width = 55

    relationships_ws = wb.create_sheet("Relationships")
    relationships_ws.append(["Type", "Column A", "Column B", "Strength", "p-value", "Significant"])
    for cell in relationships_ws[1]:
        cell.font = Font(bold=True)
    for c in result.get("correlations", []):
        relationships_ws.append(
            ["Correlation", c.get("label_a"), c.get("label_b"), c.get("r"), c.get("p_value"), c.get("significant")]
        )
    for a in result.get("associations", []):
        relationships_ws.append(
            ["Association", a.get("label_a"), a.get("label_b"), a.get("cramers_v"), a.get("p_value"), a.get("significant")]
        )
    relationships_ws.column_dimensions["B"].width = 22
    relationships_ws.column_dimensions["C"].width = 22

    alerts_ws = wb.create_sheet("Risk Alerts")
    alerts_ws.append(["Severity", "Alert"])
    for cell in alerts_ws[1]:
        cell.font = Font(bold=True)
    for a in result.get("risk_alerts", []):
        severity, message = _format_alert(a)
        alerts_ws.append([severity, message])
    alerts_ws.column_dimensions["B"].width = 70

    anomalies_ws = wb.create_sheet("Anomalies")
    anomalies_ws.append(["Column", "Row ID", "Value", "Direction"])
    for cell in anomalies_ws[1]:
        cell.font = Font(bold=True)
    for a in result.get("anomalies", []):
        anomalies_ws.append([a.get("column"), a.get("row"), a.get("value"), a.get("direction")])

    data_ws = wb.create_sheet("Raw Data")
    truncated_df = df.head(MAX_EXCEL_DATA_ROWS)
    data_ws.append(list(truncated_df.columns))
    for cell in data_ws[1]:
        cell.font = Font(bold=True)
    for row in truncated_df.itertuples(index=False):
        data_ws.append([None if pd.isna(v) else v for v in row])
    for i, col in enumerate(truncated_df.columns, start=1):
        data_ws.column_dimensions[get_column_letter(i)].width = 16

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PDF report
# ---------------------------------------------------------------------------


def _pdf_footer(filename: str) -> Any:
    def draw(canvas, doc):
        canvas.saveState()
        canvas.setFont("Helvetica", 8)
        canvas.setFillColor(colors.HexColor(MUTED_HEX))
        canvas.drawString(0.6 * inch, 0.4 * inch, f"IntelliVerse · {filename or 'Untitled dataset'}")
        canvas.drawRightString(letter[0] - 0.6 * inch, 0.4 * inch, f"Page {canvas.getPageNumber()}")
        canvas.restoreState()

    return draw


def _pdf_section_table(rows: list[list[str]], col_widths: list[int], header: bool = True) -> Table:
    table = Table(rows, colWidths=col_widths)
    style = [
        ("FONTSIZE", (0, 0), (-1, -1), 9),
        ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#e2e8f0")),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ("TOPPADDING", (0, 0), (-1, -1), 5),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("ROWBACKGROUNDS", (0, 1 if header else 0), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
    ]
    if header:
        style += [
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_HEX)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ]
    table.setStyle(TableStyle(style))
    return table


def build_pdf_report(result: dict[str, Any]) -> bytes:
    buf = io.BytesIO()
    filename = result.get("filename", "Untitled dataset")
    doc = SimpleDocTemplate(
        buf, pagesize=letter, topMargin=0.6 * inch, bottomMargin=0.7 * inch, leftMargin=0.6 * inch, rightMargin=0.6 * inch
    )
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="Muted", parent=styles["BodyText"], textColor=colors.HexColor(MUTED_HEX), fontSize=9))
    styles.add(ParagraphStyle(name="SectionHeading", parent=styles["Heading2"], textColor=colors.HexColor(BRAND_HEX), spaceBefore=14))
    story: list[Any] = []

    generated = datetime.now(timezone.utc).strftime("%B %d, %Y")

    story.append(Paragraph("IntelliVerse Analysis Report", styles["Title"]))
    story.append(Paragraph(filename, styles["Heading2"]))
    story.append(Paragraph(f"Generated {generated}", styles["Muted"]))
    story.append(Spacer(1, 10))

    quality = result.get("quality") or {}
    health = result.get("business_health") or {}
    meta_rows = [
        ["Domain", result.get("domain", "unknown")],
        ["Rows", str(result.get("row_count", ""))],
        ["Columns", str(result.get("column_count", ""))],
        ["Data quality score", str(quality.get("score", "n/a"))],
        ["Business health score", f"{health.get('overall', 'n/a')} / 100" if health else "n/a"],
    ]
    story.append(_pdf_section_table(meta_rows, [180, 300], header=False))
    story.append(Spacer(1, 18))

    # --- Business health -------------------------------------------------
    if health:
        story.append(Paragraph("Business health", styles["SectionHeading"]))
        story.append(Paragraph(
            "One score summarizing data quality, growth, forecast confidence, and risk — "
            "computed directly from the numbers in this report, no AI involved.",
            styles["Muted"],
        ))
        story.append(Spacer(1, 4))
        overall = health.get("overall", 0)
        story.append(Paragraph(
            f"<font color='{_health_color_hex(overall)}' size=20><b>{overall} / 100</b></font>", styles["BodyText"]
        ))
        component_rows = [["Component", "Score"]]
        for key, label in _COMPONENT_LABELS.items():
            value = (health.get("components") or {}).get(key)
            if value is not None:
                component_rows.append([label, str(value)])
        story.append(_pdf_section_table(component_rows, [220, 100]))
        story.append(Spacer(1, 8))

    # --- Data quality ------------------------------------------------------
    story.append(Paragraph("Data quality", styles["SectionHeading"]))
    story.append(Paragraph(
        f"Duplicate rows: {quality.get('duplicate_row_count', 0)} "
        f"({quality.get('duplicate_row_pct', 0)}% of all rows).",
        styles["BodyText"],
    ))
    invalid_values = quality.get("invalid_values") or []
    if invalid_values:
        iv_rows = [["Column", "Issue", "Rows affected"]]
        for iv in invalid_values[:10]:
            iv_rows.append([iv.get("semantic_label", iv.get("column", "")), iv.get("issue", ""), str(iv.get("count", ""))])
        story.append(Spacer(1, 4))
        story.append(_pdf_section_table(iv_rows, [110, 320, 70]))
    recommendations = quality.get("recommendations") or []
    if recommendations:
        story.append(Spacer(1, 6))
        story.append(Paragraph("<b>Recommendations</b>", styles["BodyText"]))
        for rec in recommendations[:6]:
            severity = (rec.get("severity") or "low").upper()
            sev_color = CRITICAL_HEX if severity == "HIGH" else WARNING_HEX if severity == "MEDIUM" else MUTED_HEX
            story.append(Paragraph(
                f"<font color='{sev_color}'><b>[{severity}]</b></font> {rec.get('issue', '')} — {rec.get('recommendation', '')}",
                styles["BodyText"],
            ))
    story.append(Spacer(1, 8))

    # --- Key findings ------------------------------------------------------
    story.append(Paragraph("Key findings", styles["SectionHeading"]))
    findings = result.get("ranked_findings", [])
    if findings:
        for f in findings[:12]:
            headline = f.get("headline", "")
            story.append(Paragraph(f"<b>&bull; {_plain_headline(headline)}</b>", styles["BodyText"]))
            story.append(Paragraph(f"<font color='{MUTED_HEX}' size=8>{headline}</font>", styles["BodyText"]))
    else:
        story.append(Paragraph("No ranked findings for this dataset.", styles["BodyText"]))
    story.append(Spacer(1, 8))

    # --- Root cause ----------------------------------------------------
    root_cause = result.get("root_cause")
    if root_cause and root_cause.get("dimensions"):
        story.append(Paragraph("Root cause breakdown", styles["SectionHeading"]))
        story.append(Paragraph(
            f"Which factors are most associated with variation in {root_cause.get('metric_label', 'the primary metric')}. "
            f"{root_cause.get('note', '')}",
            styles["Muted"],
        ))
        story.append(Spacer(1, 4))
        rc_rows = [["Dimension", "Variance explained", "Top segment", "Significant"]]
        for d in root_cause["dimensions"]:
            rc_rows.append([
                d.get("dimension_label", ""),
                f"{d.get('variance_explained_pct', 0)}%",
                str(d.get("top_segment", "")),
                "Yes" if d.get("significant") else "No",
            ])
        story.append(_pdf_section_table(rc_rows, [140, 110, 130, 80]))
        story.append(Spacer(1, 8))

    # --- Relationships ---------------------------------------------------
    correlations = result.get("correlations", [])
    associations = result.get("associations", [])
    if correlations or associations:
        story.append(Paragraph("Relationships", styles["SectionHeading"]))
        if correlations:
            story.append(Paragraph("<b>Numeric correlations</b>", styles["BodyText"]))
            corr_rows = [["Column A", "Column B", "r", "Strength", "Significant"]]
            for c in correlations[:8]:
                corr_rows.append([
                    c.get("label_a", ""), c.get("label_b", ""), str(c.get("r", "")),
                    c.get("strength", ""), "Yes" if c.get("significant") else "No",
                ])
            story.append(_pdf_section_table(corr_rows, [120, 120, 60, 80, 80]))
            story.append(Spacer(1, 6))
        if associations:
            story.append(Paragraph("<b>Categorical associations</b>", styles["BodyText"]))
            assoc_rows = [["Column A", "Column B", "Cramer's V", "Strength", "Significant"]]
            for a in associations[:8]:
                assoc_rows.append([
                    a.get("label_a", ""), a.get("label_b", ""), str(a.get("cramers_v", "")),
                    a.get("strength", ""), "Yes" if a.get("significant") else "No",
                ])
            story.append(_pdf_section_table(assoc_rows, [120, 120, 60, 80, 80]))
        story.append(Spacer(1, 8))

    # --- Segmentation ------------------------------------------------------
    clustering = result.get("clustering")
    if clustering and clustering.get("clusters"):
        story.append(Paragraph("Segmentation", styles["SectionHeading"]))
        story.append(Paragraph(
            f"{clustering.get('k', '?')} natural segments found (silhouette score: {clustering.get('silhouette_score', 'n/a')}).",
            styles["Muted"],
        ))
        story.append(Spacer(1, 4))
        seg_rows = [["Segment", "Size", "Profile"]]
        for c in clustering["clusters"]:
            profile = ", ".join(f"{k}: {v}" for k, v in (c.get("profile") or {}).items())
            seg_rows.append([f"Segment {c.get('cluster_id', '')}", str(c.get("size", "")), profile])
        story.append(_pdf_section_table(seg_rows, [80, 60, 380]))
        story.append(Spacer(1, 8))

    story.append(PageBreak())

    # --- Risk alerts ---------------------------------------------------
    story.append(Paragraph("Risk alerts", styles["SectionHeading"]))
    alerts = result.get("risk_alerts", [])
    if alerts:
        for a in alerts:
            severity, message = _format_alert(a)
            color = CRITICAL_HEX if severity == "Critical" else WARNING_HEX
            story.append(Paragraph(f"<font color='{color}'><b>{severity}</b></font>: {message}", styles["BodyText"]))
    else:
        story.append(Paragraph("No risk alerts triggered.", styles["BodyText"]))
    story.append(Spacer(1, 12))

    # --- Forecast ------------------------------------------------------
    forecast = result.get("forecast")
    story.append(Paragraph("Forecast", styles["SectionHeading"]))
    if forecast:
        mape = (forecast.get("validation") or {}).get("metrics", {}).get("mape")
        mape_str = f" — backtested MAPE: {mape}%" if mape is not None else ""
        story.append(Paragraph(
            f"Model: {forecast.get('method', 'n/a')} — column: {forecast.get('column', 'n/a')} — "
            f"trend: {forecast.get('trend', 'n/a')}{mape_str}",
            styles["BodyText"],
        ))
        period_comparison = result.get("period_comparison")
        if period_comparison and period_comparison.get("delta_pct") is not None:
            story.append(Paragraph(
                f"Most recent period ({period_comparison.get('current_period')}) vs. previous "
                f"({period_comparison.get('previous_period')}): {period_comparison.get('delta_pct')}% change.",
                styles["BodyText"],
            ))
        seasonality = result.get("seasonality") or {}
        if seasonality.get("detected"):
            story.append(Paragraph(
                f"Seasonality detected (lag {seasonality.get('lag')}, autocorrelation {seasonality.get('autocorrelation')}).",
                styles["BodyText"],
            ))
    else:
        story.append(Paragraph("No eligible forecast for this dataset.", styles["BodyText"]))
    story.append(Spacer(1, 12))

    # --- Data dictionary (appendix) --------------------------------------
    schema = result.get("schema", [])
    if schema:
        story.append(Paragraph("Data dictionary", styles["SectionHeading"]))
        schema_rows = [["Column", "Type", "Inferred meaning", "Confidence"]]
        for col in schema:
            schema_rows.append([
                col.get("name", ""), col.get("type", ""), col.get("semantic_label", ""),
                f"{round((col.get('confidence') or 0) * 100)}%" if isinstance(col.get("confidence"), (int, float)) else "n/a",
            ])
        story.append(_pdf_section_table(schema_rows, [130, 80, 200, 70]))

    footer = _pdf_footer(filename)
    doc.build(story, onFirstPage=footer, onLaterPages=footer)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PowerPoint report
# ---------------------------------------------------------------------------


def _style_header_row(table, ncols: int) -> None:
    for c in range(ncols):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _set_cell(table, row: int, col: int, text: str, size: int = 12, color: RGBColor | None = None, bold: bool = False) -> None:
    cell = table.cell(row, col)
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(size)
        p.font.bold = bold
        if color is not None:
            p.font.color.rgb = color
    if row % 2 == 0 and row > 0:
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_ROW


def build_pptx_report(result: dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_only_layout = prs.slide_layouts[5]
    filename = result.get("filename", "Untitled dataset")
    generated = datetime.now(timezone.utc).strftime("%B %d, %Y")
    _slide_counter = {"n": 0}

    def _add_footer(slide) -> None:
        _slide_counter["n"] += 1
        box = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35))
        tf = box.text_frame
        tf.text = f"IntelliVerse · {filename}"
        tf.paragraphs[0].font.size = Pt(9)
        tf.paragraphs[0].font.color.rgb = MUTED
        page_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.35))
        ptf = page_box.text_frame
        ptf.text = str(_slide_counter["n"])
        ptf.paragraphs[0].font.size = Pt(9)
        ptf.paragraphs[0].font.color.rgb = MUTED
        ptf.paragraphs[0].alignment = PP_ALIGN.RIGHT

    def new_slide(title: str, subtitle: str | None = None):
        slide = prs.slides.add_slide(title_only_layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        if subtitle:
            box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(12), Inches(0.5))
            tf = box.text_frame
            tf.text = subtitle
            tf.paragraphs[0].font.size = Pt(13)
            tf.paragraphs[0].font.color.rgb = MUTED
        _add_footer(slide)
        return slide

    def empty_state(slide, text: str):
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(10), Inches(1))
        tb.text_frame.text = text
        tb.text_frame.paragraphs[0].font.size = Pt(14)
        tb.text_frame.paragraphs[0].font.color.rgb = MUTED

    def bullet_box(slide, items: list[str], top: float = 1.7, size: int = 14):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(5.2))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"•  {item}"
            p.font.size = Pt(size)
            p.space_after = Pt(8)

    # --- Title slide ---------------------------------------------------
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "IntelliVerse Analysis Report"
    title_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND
    title_slide.placeholders[1].text = f"{filename}\nGenerated {generated}"

    # --- Executive summary / dataset overview ---------------------------
    kpi_slide = new_slide("Executive summary", f"{result.get('domain', 'unknown')} dataset")
    quality = result.get("quality") or {}
    health = result.get("business_health") or {}
    kpi_rows = [
        ("Rows", str(result.get("row_count", ""))),
        ("Columns", str(result.get("column_count", ""))),
        ("Data quality score", f"{quality.get('score', 'n/a')} / 100"),
        ("Duplicate rows", f"{quality.get('duplicate_row_count', 0)} ({quality.get('duplicate_row_pct', 0)}%)"),
    ]
    table_shape = kpi_slide.shapes.add_table(len(kpi_rows), 2, Inches(0.7), Inches(1.9), Inches(6.5), Inches(2.6))
    table = table_shape.table
    table.columns[0].width = Inches(3)
    table.columns[1].width = Inches(3.5)
    for r, (label, value) in enumerate(kpi_rows):
        _set_cell(table, r, 0, label, size=15, color=MUTED)
        _set_cell(table, r, 1, value, size=15, bold=True)

    if health:
        overall = health.get("overall", 0)
        score_box = kpi_slide.shapes.add_textbox(Inches(8.2), Inches(1.9), Inches(4.3), Inches(1.2))
        tf = score_box.text_frame
        tf.text = "Business Health"
        tf.paragraphs[0].font.size = Pt(15)
        tf.paragraphs[0].font.color.rgb = MUTED
        p2 = tf.add_paragraph()
        p2.text = f"{overall} / 100"
        p2.font.size = Pt(40)
        p2.font.bold = True
        p2.font.color.rgb = _health_color(overall)

    # --- Business health breakdown --------------------------------------
    if health and health.get("components"):
        bh_slide = new_slide("Business health breakdown", "Data quality, growth, forecast confidence, and risk — no AI involved")
        chart_data = CategoryChartData()
        components = health["components"]
        chart_data.categories = [_COMPONENT_LABELS.get(k, k) for k in components]
        chart_data.add_series("Score (out of 100)", list(components.values()))
        graphic_frame = bh_slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.6), Inches(1.8), Inches(11.5), Inches(4.8), chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = False
        chart.value_axis.minimum_scale = 0
        chart.value_axis.maximum_scale = 100
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = BRAND
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = Pt(13)
        plot.data_labels.font.bold = True

    # --- Data quality ----------------------------------------------------
    dq_slide = new_slide("Data quality")
    invalid_values = quality.get("invalid_values") or []
    recommendations = quality.get("recommendations") or []
    if invalid_values or recommendations:
        items = []
        for iv in invalid_values[:6]:
            items.append(f"{iv.get('semantic_label', iv.get('column', ''))}: {iv.get('issue', '')} ({iv.get('count', 0)} rows)")
        for rec in recommendations[:6]:
            severity = (rec.get("severity") or "low").upper()
            items.append(f"[{severity}] {rec.get('issue', '')} — {rec.get('recommendation', '')}")
        bullet_box(dq_slide, items, size=15)
    else:
        empty_state(dq_slide, "No data quality issues detected — score 100/100.")

    # --- Key findings ------------------------------------------------
    findings_slide = new_slide("Key findings")
    findings = result.get("ranked_findings", [])[:8]
    if findings:
        rows = len(findings) + 1
        table_shape = findings_slide.shapes.add_table(rows, 2, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))
        table = table_shape.table
        table.columns[0].width = Inches(6.5)
        table.columns[1].width = Inches(5.8)
        table.cell(0, 0).text = "Finding"
        table.cell(0, 1).text = "Statistical detail"
        _style_header_row(table, 2)
        for i, f in enumerate(findings, start=1):
            headline = f.get("headline", "")
            _set_cell(table, i, 0, _plain_headline(headline), size=13)
            _set_cell(table, i, 1, headline, size=11, color=MUTED)
    else:
        empty_state(findings_slide, "No ranked findings for this dataset.")

    # --- Root cause ------------------------------------------------------
    root_cause = result.get("root_cause")
    if root_cause and root_cause.get("dimensions"):
        rc_slide = new_slide(
            "Root cause breakdown", f"What's most associated with variation in {root_cause.get('metric_label', 'the primary metric')}"
        )
        dims = root_cause["dimensions"]
        rows = len(dims) + 1
        table_shape = rc_slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.5))
        table = table_shape.table
        for c, header in enumerate(["Dimension", "Variance explained", "Top segment", "Significant"]):
            table.cell(0, c).text = header
        _style_header_row(table, 4)
        for i, d in enumerate(dims, start=1):
            _set_cell(table, i, 0, d.get("dimension_label", ""), size=13)
            _set_cell(table, i, 1, f"{d.get('variance_explained_pct', 0)}%", size=13)
            _set_cell(table, i, 2, str(d.get("top_segment", "")), size=13)
            sig = "Yes" if d.get("significant") else "No"
            _set_cell(table, i, 3, sig, size=13, color=SUCCESS if d.get("significant") else MUTED)

    # --- Relationships ---------------------------------------------------
    correlations = result.get("correlations", [])
    associations = result.get("associations", [])
    if correlations or associations:
        rel_slide = new_slide("Relationships", "Associations found in the data — not proof of causation")
        items = []
        for c in correlations[:6]:
            sig = "significant" if c.get("significant") else "not significant"
            items.append(f"{c.get('label_a')} ↔ {c.get('label_b')}: r = {c.get('r')} ({c.get('strength')}, {sig})")
        for a in associations[:6]:
            sig = "significant" if a.get("significant") else "not significant"
            items.append(f"{a.get('label_a')} ↔ {a.get('label_b')}: Cramer's V = {a.get('cramers_v')} ({a.get('strength')}, {sig})")
        if items:
            bullet_box(rel_slide, items, size=14)
        else:
            empty_state(rel_slide, "No relationships met the significance threshold.")

    # --- Segmentation ------------------------------------------------------
    clustering = result.get("clustering")
    if clustering and clustering.get("clusters"):
        seg_slide = new_slide(
            "Segmentation", f"{clustering.get('k', '?')} natural segments (silhouette score: {clustering.get('silhouette_score', 'n/a')})"
        )
        clusters = clustering["clusters"]
        rows = len(clusters) + 1
        table_shape = seg_slide.shapes.add_table(rows, 3, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.5))
        table = table_shape.table
        table.columns[0].width = Inches(2)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(8.8)
        for c, header in enumerate(["Segment", "Size", "Profile (average values)"]):
            table.cell(0, c).text = header
        _style_header_row(table, 3)
        for i, c in enumerate(clusters, start=1):
            profile = ", ".join(f"{k}: {v}" for k, v in (c.get("profile") or {}).items())
            _set_cell(table, i, 0, f"Segment {c.get('cluster_id', '')}", size=13)
            _set_cell(table, i, 1, str(c.get("size", "")), size=13)
            _set_cell(table, i, 2, profile, size=11, color=MUTED)

    # --- Risk alerts ------------------------------------------------
    alerts_slide = new_slide("Risk alerts")
    alerts = result.get("risk_alerts", [])
    if alerts:
        rows = len(alerts) + 1
        table_shape = alerts_slide.shapes.add_table(rows, 2, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5))
        table = table_shape.table
        table.columns[0].width = Inches(2.3)
        table.columns[1].width = Inches(10)
        table.cell(0, 0).text = "Severity"
        table.cell(0, 1).text = "Alert"
        _style_header_row(table, 2)
        for i, a in enumerate(alerts, start=1):
            severity, message = _format_alert(a)
            sev_color = CRITICAL if severity == "Critical" else WARNING
            _set_cell(table, i, 0, severity, size=14, color=sev_color, bold=True)
            _set_cell(table, i, 1, message, size=13)
    else:
        empty_state(alerts_slide, "No risk alerts triggered.")

    # --- Forecast (a real chart, not text) ------------------------------
    forecast_slide = new_slide("Forecast")
    forecast = result.get("forecast")
    history = (forecast or {}).get("history") or []
    forecast_points = (forecast or {}).get("forecast") or []
    if forecast and (history or forecast_points):
        chart_data = CategoryChartData()
        chart_data.categories = [p["period"] for p in history] + [p["period"] for p in forecast_points]
        actual_series = [round(p["value"], 2) for p in history] + [None] * len(forecast_points)
        forecast_series: list[float | None] = [None] * max(len(history) - 1, 0)
        if history:
            forecast_series.append(round(history[-1]["value"], 2))
        forecast_series += [round(p["value"], 2) for p in forecast_points]
        chart_data.add_series("Actual", actual_series)
        chart_data.add_series("Forecast", forecast_series)

        graphic_frame = forecast_slide.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.3), chart_data
        )
        chart = graphic_frame.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.series[0].format.line.color.rgb = BRAND
        chart.series[1].format.line.color.rgb = WARNING

        mape = (forecast.get("validation") or {}).get("metrics", {}).get("mape")
        mape_str = f"  ·  Backtested MAPE: {mape}%" if mape is not None else ""
        caption = forecast_slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.7))
        caption.text_frame.text = (
            f"Model: {forecast.get('method', 'n/a')}  ·  Column: {forecast.get('column', 'n/a')}  ·  "
            f"Trend: {forecast.get('trend', 'n/a')}{mape_str}"
        )
        caption.text_frame.paragraphs[0].font.size = Pt(13)
        caption.text_frame.paragraphs[0].font.color.rgb = MUTED
    else:
        empty_state(forecast_slide, "No eligible forecast for this dataset.")

    # --- Anomalies ------------------------------------------------
    anomalies_slide = new_slide("Anomalies")
    anomalies = result.get("anomalies", [])[:10]
    if anomalies:
        rows = len(anomalies) + 1
        table_shape = anomalies_slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5))
        table = table_shape.table
        for c, header in enumerate(["Column", "Row", "Value", "Direction"]):
            table.cell(0, c).text = header
        _style_header_row(table, 4)
        for i, a in enumerate(anomalies, start=1):
            _set_cell(table, i, 0, a.get("column", ""), size=13)
            _set_cell(table, i, 1, a.get("row", ""), size=13)
            _set_cell(table, i, 2, a.get("value", ""), size=13)
            direction = a.get("direction", "")
            _set_cell(table, i, 3, direction, size=13, color=CRITICAL if direction == "above" else WARNING)
    else:
        empty_state(anomalies_slide, "No anomalies detected.")

    # --- Data dictionary (appendix) --------------------------------------
    schema = result.get("schema", [])
    if schema:
        schema_slide = new_slide("Data dictionary", "Every column, its inferred meaning, and detection confidence")
        rows = len(schema) + 1
        table_shape = schema_slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.9), Inches(12.3), Inches(4.8))
        table = table_shape.table
        table.columns[0].width = Inches(3)
        table.columns[1].width = Inches(2)
        table.columns[2].width = Inches(5)
        table.columns[3].width = Inches(2.3)
        for c, header in enumerate(["Column", "Type", "Inferred meaning", "Confidence"]):
            table.cell(0, c).text = header
        _style_header_row(table, 4)
        for i, col in enumerate(schema, start=1):
            confidence = col.get("confidence")
            conf_str = f"{round(confidence * 100)}%" if isinstance(confidence, (int, float)) else "n/a"
            _set_cell(table, i, 0, col.get("name", ""), size=12)
            _set_cell(table, i, 1, col.get("type", ""), size=12, color=MUTED)
            _set_cell(table, i, 2, col.get("semantic_label", ""), size=12)
            _set_cell(table, i, 3, conf_str, size=12, color=MUTED)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
